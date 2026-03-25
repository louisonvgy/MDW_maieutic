from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Couleurs
CANARD  = RGBColor(0x01, 0x6d, 0x76)
ORANGE  = RGBColor(0xec, 0x89, 0x27)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG = RGBColor(0xf1, 0xf5, 0xf9)
GREY_TXT= RGBColor(0x1e, 0x29, 0x3b)
GREY_MID= RGBColor(0x64, 0x74, 0x8b)
CANARD_L= RGBColor(0xe6, 0xf4, 0xf5)
ORANGE_L= RGBColor(0xfe, 0xf3, 0xe2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, bold=False, color=GREY_TXT,
             align=PP_ALIGN.CENTER, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_arrow(slide, x1, y, x2, color=CANARD):
    """Horizontal arrow from x1 to x2 at height y"""
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    # Simple line + triangle via connector
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

# ── FOND ─────────────────────────────────────────────────────────────────────
bg = add_rect(slide, 0, 0, 13.33, 7.5, fill=RGBColor(0xf8,0xfa,0xfc))

# ── HEADER ───────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 1.05, fill=CANARD)
add_text(slide, 'Pipeline de Clustering Sémantique des Thèses',
         0.2, 0.08, 10, 0.55, size=22, bold=True, color=WHITE)
add_text(slide, 'paraphrase-multilingual-MiniLM-L12-v2  ·  UMAP  ·  HDBSCAN  ·  TF-IDF',
         0.2, 0.6, 13, 0.38, size=11, color=CANARD_L)

# ── BLOC helper ──────────────────────────────────────────────────────────────
def bloc(slide, x, y, w, h, header, header_col, bg_col,
         lines, sub=None):
    # shadow
    add_rect(slide, x+0.04, y+0.04, w, h, fill=RGBColor(0xcc,0xd5,0xe1))
    # main box
    add_rect(slide, x, y, w, h, fill=bg_col, line=header_col, line_w=Pt(1.2))
    # header strip
    add_rect(slide, x, y, w, 0.36, fill=header_col)
    add_text(slide, header, x+0.08, y+0.03, w-0.1, 0.3,
             size=10.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # lines
    ty = y + 0.42
    for line in lines:
        add_text(slide, line, x+0.12, ty, w-0.2, 0.28,
                 size=8.5, color=GREY_TXT, align=PP_ALIGN.LEFT)
        ty += 0.27
    if sub:
        add_text(slide, sub, x+0.08, y+h-0.32, w-0.12, 0.28,
                 size=8, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

# ── ÉTIQUETTE I/O ─────────────────────────────────────────────────────────────
def io_badge(slide, x, y, text, col):
    add_rect(slide, x, y, 1.35, 0.28, fill=col, line=col, line_w=Pt(0))
    add_text(slide, text, x, y+0.02, 1.35, 0.26,
             size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── FLÈCHE droite ─────────────────────────────────────────────────────────────
def fleche(slide, x, y):
    add_rect(slide, x, y+0.09, 0.28, 0.08, fill=CANARD)
    # triangle pointe
    from pptx.util import Inches
    tri = slide.shapes.add_shape(5,  # isoceles triangle
        Inches(x+0.22), Inches(y+0.01), Inches(0.14), Inches(0.24))
    tri.fill.solid()
    tri.fill.fore_color.rgb = CANARD
    tri.line.fill.background()

# ── INPUT badge ───────────────────────────────────────────────────────────────
add_rect(slide, 0.18, 1.22, 2.1, 0.7, fill=GREY_MID, line=GREY_MID, line_w=Pt(1))
add_text(slide, '📄  INPUT', 0.22, 1.25, 2.0, 0.3,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, '9 246 titres de thèses (FR + EN)',
         0.22, 1.52, 2.0, 0.28, size=8, color=WHITE, align=PP_ALIGN.LEFT)

# Flèche 1
add_arrow(slide, 2.28, 1.57, 2.62)

# ── BLOC 1 — EMBEDDING ────────────────────────────────────────────────────────
bloc(slide, 2.65, 1.22, 2.55, 2.15,
     '① Embedding',
     CANARD, CANARD_L,
     [
         'sentence-transformers',
         'paraphrase-multilingual',
         'MiniLM-L12-v2',
         '',
         '→ 384 valeurs flottantes',
         '→ multilingue (50+ langues)',
         '→ fine-tuné paraphrases',
     ],
     sub='normalize_embeddings=True')

io_badge(slide, 2.85, 3.2, '(9 246 × 384)', CANARD)

# Flèche 2
add_arrow(slide, 5.2, 1.57, 5.54)

# ── BLOC 2 — UMAP ─────────────────────────────────────────────────────────────
bloc(slide, 5.57, 1.22, 2.55, 2.15,
     '② UMAP',
     CANARD, CANARD_L,
     [
         '384D  →  10D  (clustering)',
         '384D  →    2D  (visu)',
         '',
         'metric = cosine',
         'n_neighbors = 15',
         'min_dist = 0.0 / 0.1',
     ],
     sub='Préserve la structure non linéaire')

io_badge(slide, 5.77, 3.2, '(9 246 × 10D)', CANARD)

# Flèche 3
add_arrow(slide, 8.12, 1.57, 8.46)

# ── BLOC 3 — HDBSCAN ──────────────────────────────────────────────────────────
bloc(slide, 8.49, 1.22, 2.55, 2.15,
     '③ HDBSCAN',
     ORANGE, ORANGE_L,
     [
         'K trouvé automatiquement',
         'metric = euclidean',
         'min_cluster_size = 50',
         'min_samples = 5',
         '',
         '→ label −1 = bruit',
     ],
     sub='eom (Excess of Mass)')

io_badge(slide, 8.69, 3.2, 'labels clusters', ORANGE)

# Flèche 4
add_arrow(slide, 11.04, 1.57, 11.38)

# ── BLOC 4 — TF-IDF ───────────────────────────────────────────────────────────
bloc(slide, 11.41, 1.22, 1.72, 2.15,
     '④ TF-IDF',
     CANARD, CANARD_L,
     [
         'Mots spécifiques',
         'par cluster',
         '',
         'ngram (1,2)',
         'top 3 → label',
     ],
     sub='Labels lisibles')

# ── SÉPARATEUR ────────────────────────────────────────────────────────────────
add_rect(slide, 0.18, 3.58, 12.97, 0.03, fill=RGBColor(0xcb,0xd5,0xe1))

# ── PARTIE BASSE : OUTPUTS ────────────────────────────────────────────────────
add_text(slide, 'OUTPUTS', 0.18, 3.7, 2, 0.3,
         size=9, bold=True, color=GREY_MID, align=PP_ALIGN.LEFT)

# Output 1
add_rect(slide, 0.18, 4.08, 3.9, 1.55, fill=WHITE,
         line=CANARD, line_w=Pt(1))
add_rect(slide, 0.18, 4.08, 3.9, 0.32, fill=CANARD)
add_text(slide, 'clusters.json', 0.28, 4.1, 3.7, 0.28,
         size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide,
    '42 clusters\nid · label · keywords · nb · cx · cy · cnu_dist\n→ Bulles interactives du dashboard',
    0.28, 4.45, 3.7, 1.1, size=8.5, color=GREY_TXT, align=PP_ALIGN.LEFT)

# Output 2
add_rect(slide, 4.38, 4.08, 4.4, 1.55, fill=WHITE,
         line=CANARD, line_w=Pt(1))
add_rect(slide, 4.38, 4.08, 4.4, 0.32, fill=CANARD)
add_text(slide, 'data_clustered.json', 4.48, 4.1, 4.2, 0.28,
         size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide,
    '9 246 thèses\ndonnées + cluster_id · cluster_label · x · y\n→ Position UMAP 2D de chaque thèse',
    4.48, 4.45, 4.2, 1.1, size=8.5, color=GREY_TXT, align=PP_ALIGN.LEFT)

# Output 3 — Visu
add_rect(slide, 9.08, 4.08, 4.07, 1.55, fill=WHITE,
         line=ORANGE, line_w=Pt(1))
add_rect(slide, 9.08, 4.08, 4.07, 0.32, fill=ORANGE)
add_text(slide, 'Visualisation dashboard', 9.18, 4.1, 3.87, 0.28,
         size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide,
    'Scatter UMAP 2D → chaque thèse = 1 point\nBulles cliquables par thématique\nFiltre CNU · établissement · année',
    9.18, 4.45, 3.87, 1.1, size=8.5, color=GREY_TXT, align=PP_ALIGN.LEFT)

# ── NOTE BAS ──────────────────────────────────────────────────────────────────
add_rect(slide, 0.18, 5.82, 12.97, 0.28, fill=CANARD_L)
add_text(slide,
    'Référence : BERTopic — Grootendorst (2022) · arXiv:2203.05794',
    0.3, 5.84, 12.7, 0.24, size=7.5, italic=True, color=GREY_MID,
    align=PP_ALIGN.LEFT)

# ── PIED ──────────────────────────────────────────────────────────────────────
add_text(slide, 'MAIEUTiC · Analyse textuelle des thèses',
         9.5, 7.1, 3.6, 0.28, size=7.5, color=GREY_MID, align=PP_ALIGN.RIGHT)

prs.save('/Users/louisonvaugoyeau/MDW_maieutic/pipeline_clustering.pptx')
print('Slide générée → pipeline_clustering.pptx')
